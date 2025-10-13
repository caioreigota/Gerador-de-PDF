from flask import Flask, request, send_file, jsonify, after_this_request
import fitz  # PyMuPDF
import tempfile
import json
import requests
import zipfile
import os
import re
import shutil
import time
from pathlib import Path
from pptx import Presentation
from PIL import Image, ImageDraw, ImageChops, ImageOps
from io import BytesIO
import subprocess
import gc
from playwright.sync_api import sync_playwright
from werkzeug.utils import secure_filename
from urllib.parse import urlparse, unquote
from PIL import ImageFont
import imgkit
from io import BytesIO

# Suporte a DOCX e DOC
try:
    from docx import Document  # python-docx
    try:
        from docx.shared import Mm
    except Exception:
        Mm = None
except Exception:
    Document = None

# Imports adicionais para formatar DOCX (se disponíveis)
try:
    from docx.shared import Pt, Cm
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement
    from docx.oxml.ns import qn
except Exception:
    Pt = None
    Cm = None
    WD_ALIGN_PARAGRAPH = None
    OxmlElement = None
    qn = None

try:
    import textract  # para .doc (requer 'antiword' no sistema)
except Exception:
    textract = None

# Conversão PDF->DOCX editável (opcional)
try:
    from pdf2docx import Converter  # editable PDF->DOCX
except Exception:
    Converter = None

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # Limite de 50MB por upload

def substituir_textos(doc, substituicoes):
    for page in doc:
        insercoes = []
        aplicar_redaction = False

        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    texto_original = span["text"]
                    for chave, novo_valor in substituicoes.items():
                        marcador = f"[{chave}]"
                        if marcador in texto_original:
                            bbox = span["bbox"]
                            tamanho = span["size"]
                            cor_int = span["color"]
                            r = (cor_int >> 16) & 255
                            g = (cor_int >> 8) & 255
                            b = cor_int & 255
                            cor = (r / 255, g / 255, b / 255)

                            # Marcar para redaction
                            page.add_redact_annot(bbox, fill=(1, 1, 1), cross_out=False)
                            aplicar_redaction = True

                            novo_texto = texto_original.replace(marcador, novo_valor)
                            insercoes.append((bbox, novo_texto, tamanho, cor))

        if aplicar_redaction:
            page.apply_redactions()

        for bbox, texto, tamanho, cor in insercoes:
            x = bbox[0]
            y = bbox[1] + tamanho * 0.8  # ajuste fino vertical
            page.insert_text(
                (x, y),
                texto,
                fontsize=tamanho,
                color=cor,
                fontname="helv",
                overlay=True,
            )

@app.route('/extrair-texto', methods=['POST'])
def extrair_texto():
    """
    Recebe um arquivo (PDF, DOC ou DOCX) via multipart/form-data (campo 'file')
    e retorna o texto extraído em JSON.
    Resposta: { "texto": "..." }
    """
    uploaded = request.files.get('file')
    if not uploaded:
        return jsonify({"error": "Envie o arquivo no campo 'file' (multipart/form-data)."}), 400

    filename = secure_filename(uploaded.filename or '')
    if not filename:
        return jsonify({"error": "Nome de arquivo inválido."}), 400

    _, ext = os.path.splitext(filename)
    ext = (ext or '').lower()

    # Salva em arquivo temporário para lidar com libs que precisam de caminho (ex: textract)
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp_path = tmp.name
        uploaded.stream.seek(0)
        shutil.copyfileobj(uploaded.stream, tmp)

    try:
        if ext == '.pdf':
            texto = _extract_text_pdf(tmp_path)
        elif ext == '.docx':
            texto = _extract_text_docx(tmp_path)
        elif ext == '.doc':
            texto = _extract_text_doc(tmp_path)
        else:
            return jsonify({"error": f"Extensão não suportada: {ext}. Use PDF, DOC ou DOCX."}), 415

        return jsonify({"texto": texto})
    except Exception as e:
        return jsonify({"error": f"Falha ao extrair texto: {str(e)}"}), 500
    finally:
        try:
            os.remove(tmp_path)
        except Exception:
            pass


def _extract_text_pdf(path: str) -> str:
    doc = fitz.open(path)
    try:
        parts = []
        for page in doc:
            parts.append(page.get_text("text"))
        return "\n".join(parts).strip()
    finally:
        doc.close()


def _extract_text_docx(path: str) -> str:
    if Document is None:
        raise RuntimeError("Suporte a DOCX requer 'python-docx' instalado.")
    d = Document(path)
    # Parágrafos
    paras = [p.text for p in d.paragraphs if p.text]
    # Tabelas
    tables_text = []
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                tables_text.append("\t".join(cells))
    parts = []
    if paras:
        parts.append("\n".join(paras))
    if tables_text:
        parts.append("\n".join(tables_text))
    return "\n\n".join(parts).strip()


def _extract_text_doc(path: str) -> str:
    # Usa textract + antiword (no SO) para .doc
    if textract is None:
        raise RuntimeError("Suporte a .doc requer 'textract' instalado e o utilitário 'antiword' no sistema.")
    content = textract.process(path)  # retorna bytes
    return content.decode('utf-8', errors='replace').strip()


def _set_paragraph_bottom_border(paragraph, size_eights: int = 4, color: str = "000000"):
    """Adiciona uma borda inferior (linha) ao parágrafo via OXML.
    size_eights: espessura em oitavos de ponto (4 = 0,5 pt)
    """
    if OxmlElement is None or qn is None:
        return  # Sem suporte OXML disponível
    p_pr = paragraph._p.get_or_add_pPr()
    p_bdr = p_pr.find(qn('w:pBdr'))
    if p_bdr is None:
        p_bdr = OxmlElement('w:pBdr')
        p_pr.append(p_bdr)
    bottom = p_bdr.find(qn('w:bottom'))
    if bottom is None:
        bottom = OxmlElement('w:bottom')
        p_bdr.append(bottom)
    bottom.set(qn('w:val'), 'single')
    bottom.set(qn('w:sz'), str(size_eights))
    bottom.set(qn('w:space'), '1')
    bottom.set(qn('w:color'), color)


def _set_doc_defaults(doc):
    """Configura tamanho de página, margens e fonte padrão conforme especificação."""
    try:
        section = doc.sections[0]
        if Mm:
            section.page_width = Mm(210)
            section.page_height = Mm(297)
            # Margens 2,54 cm
            if Cm:
                cm_254 = Cm(2.54)
                section.top_margin = cm_254
                section.bottom_margin = cm_254
                section.left_margin = cm_254
                section.right_margin = cm_254
        # Fonte padrão Calibri 11 pt
        if hasattr(doc, 'styles') and Pt is not None:
            try:
                normal = doc.styles['Normal']
                normal.font.name = 'Calibri'
                normal.font.size = Pt(11)
            except Exception:
                pass
    except Exception:
        pass


def _add_section_title(doc, text: str):
    p = doc.add_paragraph()
    run = p.add_run((text or '').upper())
    if Pt is not None:
        run.font.size = Pt(12)
    run.bold = True
    if WD_ALIGN_PARAGRAPH is not None:
        p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    if Pt is not None:
        pf = p.paragraph_format
        pf.space_before = Pt(12)
        pf.space_after = Pt(6)
    return p


def _add_text_paragraphs(doc, text: str):
    """Adiciona parágrafos simples (respeitando quebras de linha)."""
    if not text:
        return
    for part in str(text).splitlines():
        p = doc.add_paragraph(part)
        if Pt is not None:
            pf = p.paragraph_format
            pf.space_after = Pt(6)


def _add_bullet(doc, text: str):
    p = doc.add_paragraph(text)
    try:
        p.style = 'List Bullet'
    except Exception:
        pass
    if Cm is not None and Pt is not None:
        pf = p.paragraph_format
        pf.left_indent = Cm(0.75)
        pf.space_after = Pt(3)
    return p


@app.route('/gerar-curriculo-docx', methods=['POST'])
def gerar_curriculo_docx():
    """Gera um .docx de currículo a partir do JSON fornecido na especificação.

    Aceita:
    - Objeto com chave 'output', ou
    - Array contendo um objeto cuja chave é 'output', ou
    - Objeto já com os campos na raiz.
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "JSON inválido ou ausente."}), 400

    # Normaliza para o objeto de saída esperado
    payload = None
    if isinstance(data, list) and data:
        item = data[0]
        if isinstance(item, dict):
            payload = item.get('output', item)
    elif isinstance(data, dict):
        payload = data.get('output', data)

    if not isinstance(payload, dict):
        return jsonify({"error": "Estrutura do JSON não reconhecida."}), 400

    if Document is None:
        return jsonify({"error": "Suporte a DOCX requer 'python-docx' instalado."}), 500

    # Se houver template URL no header, prioriza substituição no template
    template_url = (
        request.headers.get('Template-Url')
        or request.headers.get('X-Template-Url')
        or request.headers.get('template-url')
    )
    if template_url:
        try:
            resp = requests.get(template_url, timeout=30)
            resp.raise_for_status()
        except Exception as e:
            return jsonify({"error": f"Falha ao baixar template: {str(e)}"}), 400
        try:
            doc = Document(BytesIO(resp.content))
        except Exception as e:
            return jsonify({"error": f"Falha ao abrir template DOCX: {str(e)}"}), 400
        mapping = _build_text_blocks_from_payload(payload)
        _replace_placeholders_in_doc(doc, mapping)
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
        tmp_path = tmp.name
        tmp.close()
        try:
            doc.save(tmp_path)
            return send_file(
                tmp_path,
                mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                as_attachment=True,
                download_name='curriculo_preenchido.docx'
            )
        finally:
            pass

    # Caso contrário, cria documento formatado do zero
    doc = Document()
    _set_doc_defaults(doc)

    # Cabeçalho: Nome, Contatos, Linha divisória
    nome = payload.get('nome_completo') or ''
    localizacao = payload.get('localizacao')
    contatos = payload.get('contatos') or {}
    telefone = contatos.get('telefone')
    email = contatos.get('email')
    linkedin = contatos.get('linkedin')
    sites = contatos.get('sites') or []

    # Nome
    p_nome = doc.add_paragraph()
    run_nome = p_nome.add_run(str(nome or '').strip())
    run_nome.bold = True
    if Pt is not None:
        run_nome.font.size = Pt(14)
    # sem espaço adicional
    if Pt is not None:
        pf = p_nome.paragraph_format
        pf.space_before = Pt(0)
        pf.space_after = Pt(0)

    # Linha de contato
    contato_parts = []
    if localizacao:
        contato_parts.append(str(localizacao))
    if telefone:
        contato_parts.append(str(telefone))
    if email:
        contato_parts.append(str(email))
    if linkedin:
        contato_parts.append(str(linkedin))
    if sites:
        contato_parts.extend([str(s) for s in sites if s])
    contato_text = ' | '.join(contato_parts)
    p_contato = doc.add_paragraph(contato_text)
    if Pt is not None:
        p_contato.paragraph_format.space_before = Pt(0)
        p_contato.paragraph_format.space_after = Pt(0)
        for r in p_contato.runs:
            r.font.size = Pt(10.5)

    # Linha divisória (0,5 pt, espaçamento 6 pt acima/abaixo)
    p_rule = doc.add_paragraph()
    if Pt is not None:
        p_rule.paragraph_format.space_before = Pt(6)
        p_rule.paragraph_format.space_after = Pt(6)
    _set_paragraph_bottom_border(p_rule, size_eights=4, color='000000')

    # Seções do corpo
    objetivo = payload.get('objetivo_profissional')
    resumo = payload.get('resumo_profissional')
    competencias = payload.get('competencias_chave') or []
    formacao = payload.get('formacao_academica') or []
    experiencias = payload.get('experiencia_profissional') or []
    idiomas = payload.get('idiomas') or []
    conquistas = payload.get('conquistas_certificados') or []
    adicionais = payload.get('informacoes_adicionais') or []

    if objetivo:
        _add_section_title(doc, 'Objetivo')
        _add_text_paragraphs(doc, str(objetivo))

    if resumo:
        _add_section_title(doc, 'Resumo Profissional')
        _add_text_paragraphs(doc, str(resumo))

    if competencias:
        _add_section_title(doc, 'Competências Chave')
        for comp in competencias:
            _add_bullet(doc, str(comp))

    if formacao:
        _add_section_title(doc, 'Educação')
        for item in formacao:
            curso = (item or {}).get('curso')
            instituicao = (item or {}).get('instituicao')
            periodo = (item or {}).get('periodo')
            p = doc.add_paragraph()
            # "Curso, Instituição (Período)" com instituição em negrito
            if curso:
                p.add_run(str(curso))
                p.add_run(', ')
            if instituicao:
                r_inst = p.add_run(str(instituicao))
                r_inst.bold = True
            if periodo:
                p.add_run(f" ({periodo})")
            if Pt is not None:
                p.paragraph_format.space_after = Pt(3)

    if experiencias:
        _add_section_title(doc, 'Experiência Profissional')
        for exp in experiencias:
            empresa = (exp or {}).get('empresa')
            cargo = (exp or {}).get('cargo')
            periodo = (exp or {}).get('periodo')
            resp = (exp or {}).get('responsabilidades_resultados') or []
            # Linha de cabeçalho do item: Empresa — Cargo (Período)
            p_head = doc.add_paragraph()
            if empresa:
                r_emp = p_head.add_run(str(empresa))
                r_emp.bold = True
            if cargo:
                p_head.add_run(' — ' + str(cargo))
            if periodo:
                p_head.add_run(' (' + str(periodo) + ')')
            if Pt is not None:
                p_head.paragraph_format.space_after = Pt(3)
            # Bullets das responsabilidades
            for r in resp:
                _add_bullet(doc, str(r))

    if idiomas:
        _add_section_title(doc, 'Idiomas')
        for idm in idiomas:
            idioma = (idm or {}).get('idioma')
            nivel = (idm or {}).get('nivel')
            p = doc.add_paragraph()
            parts = []
            if idioma:
                parts.append(str(idioma))
            if nivel:
                parts.append('— ' + str(nivel))
            p.add_run(' '.join(parts))
            if Pt is not None:
                p.paragraph_format.space_after = Pt(3)

    if conquistas:
        _add_section_title(doc, 'Conquistas e Certificados')
        for c in conquistas:
            _add_bullet(doc, str(c))

    if adicionais:
        _add_section_title(doc, 'Informações Adicionais')
        for a in adicionais:
            _add_bullet(doc, str(a))

    # Salva temporário e retorna
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
    tmp_path = tmp.name
    tmp.close()
    try:
        doc.save(tmp_path)
        return send_file(
            tmp_path,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True,
            download_name='curriculo.docx'
        )
    finally:
        # O arquivo temporário será removido pelo servidor/OS posteriormente; evitar remover antes de enviar
        pass


def _build_text_blocks_from_payload(payload: dict) -> dict:
    """Monta mapeamento de placeholders -> texto a partir do JSON de entrada."""
    contatos = (payload or {}).get('contatos') or {}
    # Campos simples
    nome = (payload or {}).get('nome_completo') or ''
    localizacao = (payload or {}).get('localizacao') or ''
    telefone = contatos.get('telefone') or ''
    email = contatos.get('email') or ''
    linkedin = contatos.get('linkedin') or ''
    objetivo = (payload or {}).get('objetivo_profissional') or ''
    resumo = (payload or {}).get('resumo_profissional') or ''

    # Formação
    formacao_list = (payload or {}).get('formacao_academica') or []
    formacao_lines = []
    for f in formacao_list:
        if not isinstance(f, dict):
            continue
        curso = f.get('curso')
        inst = f.get('instituicao')
        periodo = f.get('periodo')
        parts = []
        if curso:
            parts.append(str(curso))
        if inst:
            parts.append(str(inst))
        line = ', '.join(parts)
        if periodo:
            line = f"{line} ({periodo})" if line else f"({periodo})"
        if line:
            formacao_lines.append(line)
    formacao_text = "\n".join(formacao_lines)

    # Experiência
    exp_list = (payload or {}).get('experiencia_profissional') or []
    exp_blocks = []
    for e in exp_list:
        if not isinstance(e, dict):
            continue
        empresa = e.get('empresa')
        cargo = e.get('cargo')
        periodo = e.get('periodo')
        head_parts = []
        if empresa:
            head_parts.append(str(empresa))
        if cargo:
            head_parts.append('— ' + str(cargo))
        head = ' '.join(head_parts)
        if periodo:
            head = f"{head} ({periodo})" if head else f"({periodo})"
        lines = []
        if head:
            lines.append(head)
        for r in (e.get('responsabilidades_resultados') or []):
            if r is None:
                continue
            lines.append(f"• {str(r)}")
        if lines:
            exp_blocks.append("\n".join(lines))
    exp_text = "\n\n".join(exp_blocks)

    # Idiomas
    idiomas = (payload or {}).get('idiomas') or []
    idiomas_lines = []
    for i in idiomas:
        if not isinstance(i, dict):
            continue
        idioma = i.get('idioma')
        nivel = i.get('nivel')
        if idioma and nivel:
            idiomas_lines.append(f"{idioma} — {nivel}")
        elif idioma:
            idiomas_lines.append(str(idioma))
    idiomas_text = "\n".join(idiomas_lines)

    # Certificações
    certs = (payload or {}).get('conquistas_certificados') or []
    cert_lines = [f"• {str(c)}" for c in certs if c]
    certs_text = "\n".join(cert_lines)

    mapping = {
        'NOME': str(nome),
        'CIDADE': str(localizacao),
        'TELEFONE': str(telefone),
        'TELEDONE': str(telefone),  # cobre o placeholder com possível typo
        'EMAIL': str(email),
        'LINKEDIN': str(linkedin),
        'OBJETIVO': str(objetivo),
        'RESUMO_PROFISSIONAL': str(resumo),
        'FORMACAO': formacao_text,
        'EXPERIENCIA_PROFISSIONAL': exp_text,
        'IDIOMAS': idiomas_text,
        'CERTIFICACOES': certs_text,
        'CERTIFICAÇÕES': certs_text,
    }
    # Remove chaves com valor vazio para evitar inserir 'None' ou espaços supérfluos
    return {k: v for k, v in mapping.items() if v is not None}


def _replace_placeholders_in_paragraph(paragraph, mapping_upper: dict):
    # Concatena runs
    full = ''.join(run.text for run in paragraph.runs) if paragraph.runs else paragraph.text
    if not full:
        return
    changed = full
    for key, val in mapping_upper.items():
        placeholder = f"[{key}]"
        if placeholder in changed:
            changed = changed.replace(placeholder, val)
    if changed != full:
        # Atualiza runs de forma simplificada: 1 run com o texto final
        if paragraph.runs:
            paragraph.runs[0].text = changed
            for r in paragraph.runs[1:]:
                r.text = ''
        else:
            paragraph.text = changed


def _replace_placeholders_in_doc(doc, mapping: dict):
    # Normaliza chaves para maiúsculas
    mp = {str(k).upper(): str(v) for k, v in mapping.items() if v is not None}
    # Parágrafos do corpo
    for p in doc.paragraphs:
        _replace_placeholders_in_paragraph(p, mp)
    # Dentro de tabelas, se houver
    for t in doc.tables:
        for row in t.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    _replace_placeholders_in_paragraph(p, mp)


@app.route('/gerar-docx-de-template', methods=['POST'])
def gerar_docx_de_template():
    """Baixa um arquivo .docx de 'template_url' e substitui placeholders [CHAVE]
    com valores derivados do JSON de currículo.

    Body JSON:
    {
      "template_url": "https://.../modelo.docx",
      "output": { ... mesmo schema ... }
    }
    Também aceita o JSON direto na raiz (sem 'output') ou dentro de uma lista.
    """
    data = request.get_json(silent=True)
    if not data:
        return jsonify({"error": "JSON inválido ou ausente."}), 400

    template_url = data.get('template_url') if isinstance(data, dict) else None
    if not template_url:
        return jsonify({"error": "Campo 'template_url' é obrigatório."}), 400

    # Extrai payload como antes
    payload = None
    if isinstance(data, list) and data:
        item = data[0]
        if isinstance(item, dict):
            payload = item.get('output', item)
    elif isinstance(data, dict):
        payload = data.get('output', data)
    if not isinstance(payload, dict):
        return jsonify({"error": "Estrutura do JSON de dados não reconhecida."}), 400

    if Document is None:
        return jsonify({"error": "Suporte a DOCX requer 'python-docx' instalado."}), 500

    try:
        resp = requests.get(template_url, timeout=30)
        resp.raise_for_status()
    except Exception as e:
        return jsonify({"error": f"Falha ao baixar template: {str(e)}"}), 400

    try:
        doc = Document(BytesIO(resp.content))
    except Exception as e:
        return jsonify({"error": f"Falha ao abrir template DOCX: {str(e)}"}), 400

    mapping = _build_text_blocks_from_payload(payload)
    _replace_placeholders_in_doc(doc, mapping)

    tmp = tempfile.NamedTemporaryFile(delete=False, suffix='.docx')
    tmp_path = tmp.name
    tmp.close()
    try:
        doc.save(tmp_path)
        return send_file(
            tmp_path,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True,
            download_name='curriculo_preenchido.docx'
        )
    finally:
        pass

@app.route('/pdf-para-imagem', methods=['POST'])
def pdf_para_imagem():
    data = request.get_json()
    if not data or 'pdf_url' not in data:
        return {'error': 'pdf_url é obrigatório'}, 400

    try:
        response = requests.get(data['pdf_url'])
        response.raise_for_status()
    except Exception as e:
        return {'error': f'Erro ao baixar PDF: {str(e)}'}, 400

    try:
        parsed_url = urlparse(data['pdf_url'])
        pdf_filename = os.path.basename(parsed_url.path)
        pdf_filename = unquote(pdf_filename)
        nome_base = os.path.splitext(pdf_filename)[0]
        nome_imagem = f"{nome_base}.png"

        doc = fitz.open(stream=response.content, filetype="pdf")
        if len(doc) == 0:
            return {'error': 'PDF sem páginas'}, 400

        page = doc[0]
        pix = page.get_pixmap(dpi=150)

        img_bytes = BytesIO()
        img_bytes.write(pix.tobytes("png"))
        img_bytes.seek(0)

        doc.close()
        gc.collect()

        return send_file(
            img_bytes,
            mimetype="image/png",
            as_attachment=True,
            download_name=nome_imagem
        )

    except Exception as e:
        return {'error': f'Erro ao processar PDF: {str(e)}'}, 500


@app.route('/preencher-pdf-url', methods=['POST'])
def preencher_pdf_url():
    data = request.get_json()
    if not data or 'pdf_url' not in data or 'substituicoes' not in data:
        return {'error': 'pdf_url e substituicoes são obrigatórios'}, 400

    try:
        response = requests.get(data['pdf_url'])
        response.raise_for_status()
    except Exception as e:
        return {'error': f'Erro ao baixar PDF: {str(e)}'}, 400

    substituicoes = data['substituicoes']
    if not isinstance(substituicoes, dict) or not substituicoes:
        return {'error': 'Substituições inválidas ou vazias'}, 400

    try:
        doc = fitz.open(stream=response.content, filetype="pdf")
        substituir_textos(doc, substituicoes)

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_out:
            doc.save(tmp_out.name, deflate=True, garbage=4, clean=True)
            doc.close()
            gc.collect()
            return send_file(
                tmp_out.name,
                mimetype="application/pdf",
                as_attachment=True,
                download_name="preenchido.pdf"
            )
    except Exception as e:
        return {'error': f'Erro ao processar PDF: {str(e)}'}, 500
    

@app.route('/preencher-html-url', methods=['POST'])
def preencher_html_url():
    data = request.get_json()
    if not data or 'html_url' not in data or 'substituicoes' not in data:
        return {'error': 'html_url e substituicoes são obrigatórios'}, 400

    try:
        response = requests.get(data['html_url'])
        response.raise_for_status()
        html = response.text
    except Exception as e:
        return {'error': f'Erro ao baixar HTML: {str(e)}'}, 400

    substituicoes = data['substituicoes']
    if not isinstance(substituicoes, dict) or not substituicoes:
        return {'error': 'Substituições inválidas ou vazias'}, 400

    try:
        for chave, valor in substituicoes.items():
            html = html.replace(f"{{{{ {chave} }}}}", valor)

        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as img_file:
            imgkit.from_string(html, img_file.name, options={"format": "png", "width": 1080, "height": 1080})
            return send_file(
                img_file.name,
                mimetype="image/png",
                as_attachment=True,
                download_name="imagem_gerada.png"
            )
    except Exception as e:
        return {'error': f'Erro ao processar HTML: {str(e)}'}, 500

@app.route('/gerar-imagem-vaga', methods=['POST'])
def gerar_imagem_vaga():
    def draw_text_wrap(draw, text, font, max_width, x, y, line_spacing=10):
        words = text.split()
        line = ""
        lines = []
        for word in words:
            test_line = f"{line} {word}".strip()
            width, _ = draw.textsize(test_line, font=font)
            if width <= max_width:
                line = test_line
            else:
                lines.append(line)
                line = word
        lines.append(line)

        for i, line in enumerate(lines):
            draw.text((x, y + i * (font.size + line_spacing)), line, font=font, fill="black")

    data = request.get_json()
    if not data or 'pdf_url' not in data or 'substituicoes' not in data:
        return {'error': 'pdf_url e substituicoes são obrigatórios'}, 400

    try:
        response = requests.get(data['pdf_url'])
        response.raise_for_status()
    except Exception as e:
        return {'error': f'Erro ao baixar PDF: {str(e)}'}, 400

    try:
        doc = fitz.open(stream=response.content, filetype="pdf")
        page = doc[0]
        pix = page.get_pixmap(dpi=300)
        img = Image.open(BytesIO(pix.tobytes("png"))).convert("RGB")
        draw = ImageDraw.Draw(img)

        # Fonte segura para ambientes Linux/Docker
        try:
            font_padrao = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", size=28)
        except:
            font_padrao = ImageFont.load_default()

        campos = data["substituicoes"]

        # Coordenadas ajustadas com base no modelo
        coordenadas = {
            "cargo": (130, 420),
            "complemento": (130, 470),
            "Requisito 1": (130, 540),
            "Requisito 2": (130, 580),
            "Requisito 3": (130, 620),
            "Requisito 4": (130, 660),
            "Requisito 5": (130, 700),
            "localizacao": (200, 820),
            "modalidade": (530, 820)
        }

        for chave, pos in coordenadas.items():
            texto = campos.get(chave, "")
            if texto:
                draw_text_wrap(draw, texto, font_padrao, max_width=580, x=pos[0], y=pos[1])

        img_bytes = BytesIO()
        img.save(img_bytes, format="PNG")
        img_bytes.seek(0)
        doc.close()
        gc.collect()

        return send_file(
            img_bytes,
            mimetype="image/png",
            as_attachment=True,
            download_name="vaga_preenchida.png"
        )

    except Exception as e:
        return {'error': f'Erro ao gerar imagem: {str(e)}'}, 500

@app.route('/pptx-para-imagens', methods=['POST'])
def pptx_para_imagens():
    data = request.get_json()
    if not data or 'pptx_url' not in data:
        return {'error': 'pptx_url é obrigatório'}, 400

    try:
        response = requests.get(data['pptx_url'])
        response.raise_for_status()
    except Exception as e:
        return {'error': f'Erro ao baixar PPTX: {str(e)}'}, 400

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_pptx:
            tmp_pptx.write(response.content)
            pptx_path = tmp_pptx.name

        output_dir = tempfile.mkdtemp()

        subprocess.run([
            "libreoffice",
            "--headless",
            "--convert-to", "png",
            "--outdir", output_dir,
            pptx_path
        ], check=True)

        zip_path = os.path.join(output_dir, "slides.zip")
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for file_name in sorted(os.listdir(output_dir)):
                if file_name.endswith(".png"):
                    file_path = os.path.join(output_dir, file_name)
                    zipf.write(file_path, arcname=file_name)

        gc.collect()

        return send_file(
            zip_path,
            mimetype="application/zip",
            as_attachment=True,
            download_name="slides_imagens.zip"
        )

    except subprocess.CalledProcessError as e:
        return {'error': f'Erro ao converter com LibreOffice: {str(e)}'}, 500
    except Exception as e:
        return {'error': f'Erro ao processar PPTX: {str(e)}'}, 500


@app.route('/redimensionar-imagem', methods=['POST'])
def redimensionar_imagem():
    if 'image' not in request.files:
        return {'error': 'Imagem é obrigatória'}, 400

    largura = request.form.get('largura') or request.form.get('width')
    altura = request.form.get('altura') or request.form.get('height')
    if not largura or not altura:
        return {'error': 'Largura e altura são obrigatórios'}, 400

    try:
        largura = int(largura)
        altura = int(altura)
        if largura <= 0 or altura <= 0:
            raise ValueError
    except ValueError:
        return {'error': 'Largura e altura devem ser números inteiros positivos'}, 400

    try:
        img_file = request.files['image']
        img = Image.open(img_file.stream)
        img = img.resize((largura, altura))

        output = BytesIO()
        formato = img.format or 'PNG'
        img.save(output, format=formato)
        output.seek(0)

        mimetype = img_file.mimetype or f'image/{formato.lower()}'
        ext = formato.lower()
        return send_file(
            output,
            mimetype=mimetype,
            as_attachment=True,
            download_name=f'resized.{ext}'
        )
    except Exception as e:
        return {'error': f'Erro ao redimensionar imagem: {str(e)}'}, 500


@app.route('/cortar-redimensionar-imagem', methods=['POST'])
def cortar_redimensionar_imagem():
    if 'image' not in request.files:
        return {'error': 'Imagem é obrigatória'}, 400

    largura = request.form.get('largura') or request.form.get('width')
    altura = request.form.get('altura') or request.form.get('height')
    top = request.form.get('top') or request.form.get('crop_top') or '0'
    right = request.form.get('right') or request.form.get('crop_right') or '0'
    left = request.form.get('left') or request.form.get('crop_left') or '0'
    bottom = (
        request.form.get('bottom')
        or request.form.get('footer')
        or request.form.get('crop_bottom')
        or '0'
    )
    if not largura or not altura:
        return {'error': 'Largura e altura são obrigatórios'}, 400

    try:
        largura = int(largura)
        altura = int(altura)
        top = int(top)
        right = int(right)
        left = int(left)
        bottom = int(bottom)
        if (
            largura <= 0
            or altura <= 0
            or top < 0
            or right < 0
            or left < 0
            or bottom < 0
        ):
            raise ValueError
    except ValueError:
        return {'error': 'Parâmetros de corte devem ser inteiros não negativos e largura/altura positivos'}, 400

    try:
        img_file = request.files['image']
        img = Image.open(img_file.stream).convert('RGB')
        background_color = img.getpixel((0, 0))

        def crop_blank_top_bottom(im, threshold=250, blank_ratio=0.99):
            gray = im.convert('L')
            width, height = gray.size
            pixels = gray.load()

            def row_is_blank(y):
                blank = 0
                for x in range(width):
                    if pixels[x, y] >= threshold:
                        blank += 1
                return blank / width >= blank_ratio

            top = 0
            bottom = height - 1
            while top < height and row_is_blank(top):
                top += 1
            while bottom > top and row_is_blank(bottom):
                bottom -= 1
            return im.crop((0, top, width, bottom + 1))

        img = crop_blank_top_bottom(img)


        width, height = img.size
        left = max(0, min(left, width - 1))
        top = max(0, min(top, height - 1))
        right = max(0, min(right, width - left - 1))
        bottom = max(0, min(bottom, height - top - 1))
        img = img.crop((left, top, width - right, height - bottom))


        img.thumbnail((largura, altura), Image.LANCZOS)

        background = Image.new("RGB", (largura, altura), background_color)
        offset = ((largura - img.size[0]) // 2, (altura - img.size[1]) // 2)
        background.paste(img, offset)

        output = BytesIO()
        formato = img_file.mimetype.split('/')[-1].upper() if img_file.mimetype else 'PNG'
        background.save(output, format=formato)
        output.seek(0)

        mimetype = img_file.mimetype or f'image/{formato.lower()}'
        ext = formato.lower()
        return send_file(
            output,
            mimetype=mimetype,
            as_attachment=True,
            download_name=f'resized.{ext}'
        )
    except Exception as e:
        return {'error': f'Erro ao processar imagem: {str(e)}'}, 500


@app.route('/pptx-para-pdf', methods=['POST'])
def pptx_para_pdf():
    data = request.get_json()
    if not data or 'pptx_url' not in data or 'substituicoes' not in data:
        return {'error': 'pptx_url e substituicoes são obrigatórios'}, 400

    try:
        response = requests.get(data['pptx_url'])
        response.raise_for_status()
    except Exception as e:
        return {'error': f'Erro ao baixar PPTX: {str(e)}'}, 400

    try:
        substituicoes = data['substituicoes']
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_pptx:
            tmp_pptx.write(response.content)
            pptx_path = tmp_pptx.name

        prs = Presentation(pptx_path)
        imagens = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for par in shape.text_frame.paragraphs:
                        for run in par.runs:
                            for chave, valor in substituicoes.items():
                                marcador = f"[{chave}]"
                                if marcador in run.text:
                                    run.text = run.text.replace(marcador, valor)

        for slide in prs.slides:
            img = Image.new("RGB", (1280, 720), color="white")
            draw = ImageDraw.Draw(img)
            y = 20
            for shape in slide.shapes:
                if shape.has_text_frame:
                    draw.text((20, y), shape.text, fill="black")
                    y += 30
            img_bytes = BytesIO()
            img.save(img_bytes, format="PNG", optimize=True)
            img_bytes.seek(0)
            imagens.append(img_bytes)

        pdf_path = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf").name
        doc = fitz.open()
        for img_bytes in imagens:
            img = Image.open(img_bytes)
            rect = fitz.Rect(0, 0, img.width, img.height)
            page = doc.new_page(width=img.width, height=img.height)
            page.insert_image(rect, stream=img_bytes.read())
        doc.save(pdf_path)
        doc.close()
        gc.collect()

        return send_file(
            pdf_path,
            mimetype="application/pdf",
            as_attachment=True,
            download_name="slides_convertidos.pdf"
        )

    except Exception as e:
        return {'error': f'Erro ao processar PPTX: {str(e)}'}, 500


# ================= Vídeo a partir de HTML =================

# =============== Defaults (pode ajustar) ===============
TARGET_FPS       = 30        # fps do MP4 final
CONTENT_SECONDS  = 10.0      # conteúdo útil aproximado
BUFFER_HEAD_S    = 2.0       # buffer antes
BUFFER_TAIL_S    = 2.0       # buffer depois
SCENE_THRESHOLD  = 0.0015    # sensibilidade do auto-trim do início
HEAD_PAD_S       = 0.05      # margem antes do 1º movimento
TAIL_PAD_S       = 0.10      # margem após o último movimento (apenas se cortar o final)
AUTO_TRIM_HEAD   = True      # corta automaticamente o início ESTÁTICO
AUTO_TRIM_TAIL   = False     # não corta o final por padrão
ZERO_ANIM_DELAY  = True      # zera animation/transition-delay (padrão; pode sobrescrever por requisição)
WAIT_NETWORK_IDLE= True
TIMEZONE_ID      = "America/Sao_Paulo"
AUTO_SIZE_BODY   = True      # mede o <body> e grava no tamanho exato
MAX_DIM          = 4096      # limite por dimensão para o viewport gravado
# Qualidade/compat do MP4
CRF              = 20
PRESET           = "medium"
PIX_FMT          = "yuv420p"
PROFILE          = "high"
LEVEL            = "4.0"
# =======================================================


def ensure_ffmpeg():
    if shutil.which("ffmpeg") is None or shutil.which("ffprobe") is None:
        raise RuntimeError("FFmpeg/ffprobe não encontrados no PATH.")


def file_url(p: Path) -> str:
    return p.resolve().as_uri()


def has_wkhtmltopdf() -> bool:
    return shutil.which("wkhtmltopdf") is not None


def has_wkhtmltoimage() -> bool:
    """Verifica se o binário wkhtmltoimage está disponível no PATH."""
    return shutil.which("wkhtmltoimage") is not None


def measure_body_size(html_path: Path) -> tuple[int, int]:
    """Abre SEM gravação, mede o tamanho real do <body> e devolve (w,h)."""
    with sync_playwright() as p:
        try:
            browser = p.chromium.launch(headless=True)
        except Exception:
            browser = p.chromium.launch(headless=True, args=['--no-sandbox'])  # Docker fallback
        ctx = browser.new_context(
            viewport={"width": 1080, "height": 1350},
            device_scale_factor=1.0,
            java_script_enabled=True,
            timezone_id=TIMEZONE_ID,
        )
        page = ctx.new_page()
        page.goto(file_url(html_path), wait_until="load")
        if WAIT_NETWORK_IDLE:
            try:
                page.wait_for_load_state("networkidle", timeout=8000)
            except Exception:
                pass
        try:
            page.evaluate("document.fonts && document.fonts.ready && document.fonts.ready.then(()=>{})")
            page.wait_for_timeout(50)
        except Exception:
            pass
        dims = page.evaluate(
            """
        () => {
          const b = document.body, d = document.documentElement;
          const w = Math.max(b.scrollWidth, d.scrollWidth, b.offsetWidth, d.offsetWidth, d.clientWidth);
          const h = Math.max(b.scrollHeight, d.scrollHeight, b.offsetHeight, d.offsetHeight, d.clientHeight);
          return {w, h};
        }
        """
        )
        ctx.close()
        browser.close()
    w = max(1, int(dims["w"]))
    h = max(1, int(dims["h"]))
    if w > MAX_DIM or h > MAX_DIM:
        r = min(MAX_DIM / w, MAX_DIM / h)
        w, h = max(1, int(w * r)), max(1, int(h * r))
    return w, h


def prepare_page(page, *, zero_anim_delay: bool, extra_wait_ms: int = 0):
    page.wait_for_load_state("load")
    if WAIT_NETWORK_IDLE:
        try:
            page.wait_for_load_state("networkidle", timeout=8000)
        except Exception:
            pass
    # Aguarda carregamento de fontes
    try:
        page.wait_for_function(
            "document.fonts && document.fonts.status === 'loaded'",
            timeout=8000,
        )
    except Exception:
        pass
    # Aguarda todas as imagens ficarem completas
    try:
        page.wait_for_function(
            "Array.from(document.images || []).every(img => img.complete)",
            timeout=8000,
        )
    except Exception:
        pass
    # Espera estabilidade de layout (scrollWidth/scrollHeight estáveis por alguns frames)
    try:
        page.wait_for_function(
            """
            (() => {
              if (!window.__layoutStable) { window.__layoutStable = { last: '', count: 0 }; }
              const w = document.documentElement.scrollWidth || document.body.scrollWidth;
              const h = document.documentElement.scrollHeight || document.body.scrollHeight;
              const cur = w + 'x' + h;
              if (window.__layoutStable.last === cur) { window.__layoutStable.count++; } else { window.__layoutStable.last = cur; window.__layoutStable.count = 0; }
              return window.__layoutStable.count > 10;
            })()
            """,
            timeout=3000,
        )
    except Exception:
        pass
    if zero_anim_delay:
        page.add_style_tag(
            content="""
            * { animation-delay: 0s !important; transition-delay: 0s !important; }
        """
        )
    if extra_wait_ms:
        try:
            page.wait_for_timeout(int(extra_wait_ms))
        except Exception:
            pass


def chromium_screenshot(
    html_path: Path,
    *,
    width: int | None,
    height: int | None,
    transparent: bool,
    dpr: float = 1.0,
    out_dir: Path | None = None,
    full_page: bool = False,
    media: str | None = None,
    css_inject: str | None = None,
) -> Path:
    """Renderiza o HTML com Chromium/Playwright e retorna o caminho de um PNG temporário."""
    # Determina tamanho de renderização
    if AUTO_SIZE_BODY and (width is None or height is None):
        w, h = measure_body_size(html_path)
    else:
        w = width or 1200
        h = height or 1350

    if out_dir is None:
        fd, pth = tempfile.mkstemp(suffix=".png")
        os.close(fd)
        tmp_png = Path(pth)
    else:
        tmp_png = Path(out_dir) / "_base_playwright.png"
    with sync_playwright() as p:
        try:
            browser = p.chromium.launch(headless=True)
        except Exception:
            browser = p.chromium.launch(headless=True, args=['--no-sandbox'])  # Docker fallback
        ctx = browser.new_context(
            viewport={"width": w, "height": h},
            device_scale_factor=dpr,
            java_script_enabled=True,
            timezone_id=TIMEZONE_ID,
        )
        page = ctx.new_page()
        page.goto(file_url(html_path), wait_until="load")
        prepare_page(page, zero_anim_delay=True)
        # Injeta CSS de fallback de fontes/ajustes, se fornecido
        if css_inject:
            try:
                page.add_style_tag(content=css_inject)
            except Exception:
                pass
        # Emula mídia para aplicar @media corretamente (padrão: screen)
        if media not in ("screen", "print"):
            media = "screen"
        try:
            page.emulate_media(media=media)
        except Exception:
            pass
            try:
                page.emulate_media(media=media)
            except Exception:
                pass
        # Ajuste dinamicamente se mediu body
        if AUTO_SIZE_BODY and (width is None or height is None):
            dims = page.evaluate(
                """
                () => {
                  const b = document.body, d = document.documentElement;
                  const w = Math.max(b.scrollWidth, d.scrollWidth, b.offsetWidth, d.offsetWidth, d.clientWidth);
                  const h = Math.max(b.scrollHeight, d.scrollHeight, b.offsetHeight, d.offsetHeight, d.clientHeight);
                  return {w: Math.max(1, Math.floor(w)), h: Math.max(1, Math.floor(h))};
                }
                """
            )
            w = min(MAX_DIM, int(dims.get('w') or w))
            h = min(MAX_DIM, int(dims.get('h') or h))
            page.set_viewport_size({"width": w, "height": h})
        page.screenshot(path=str(tmp_png), full_page=bool(full_page), omit_background=bool(transparent))
        ctx.close(); browser.close()
    return tmp_png


def record_webm(
    html_path: Path,
    total_seconds: float,
    width: int,
    height: int,
    out_dir: Path,
    *,
    zero_anim_delay: bool,
) -> Path:
    with sync_playwright() as p:
        try:
            browser = p.chromium.launch(headless=True)
        except Exception:
            browser = p.chromium.launch(headless=True, args=['--no-sandbox'])  # Docker fallback
        ctx = browser.new_context(
            viewport={"width": width, "height": height},
            record_video_dir=str(out_dir),
            record_video_size={"width": width, "height": height},
            device_scale_factor=1.0,
            java_script_enabled=True,
            timezone_id=TIMEZONE_ID,
        )
        page = ctx.new_page()
        page.goto(file_url(html_path), wait_until="load")
        prepare_page(page, zero_anim_delay=zero_anim_delay)
        page.wait_for_timeout(500)  # warmup
        page.wait_for_timeout(int(total_seconds * 1000))
        page.close()
        ctx.close()
        browser.close()
    vids = sorted(out_dir.rglob("*.webm"), key=lambda p: p.stat().st_mtime, reverse=True)
    if not vids:
        raise RuntimeError("Nenhum WEBM gravado.")
    return vids[0]


def video_duration(path: Path) -> float:
    out = subprocess.check_output(
        [
            "ffprobe",
            "-v",
            "error",
            "-show_entries",
            "format=duration",
            "-of",
            "default=nokey=1:noprint_wrappers=1",
            str(path),
        ]
    ).decode().strip()
    try:
        return float(out)
    except Exception:
        return 0.0


def detect_scene_changes(path: Path, threshold: float) -> list[float]:
    cmd = [
        "ffmpeg",
        "-hide_banner",
        "-i",
        str(path),
        "-filter:v",
        f"select='gt(scene,{threshold})',showinfo",
        "-an",
        "-f",
        "null",
        "-",
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)
    times = []
    for line in proc.stderr.splitlines():
        m = re.search(r"pts_time:([0-9]+\.[0-9]+)", line)
        if m:
            try:
                times.append(float(m.group(1)))
            except Exception:
                pass
    return times


def compute_trim(
    webm: Path,
    *,
    auto_head=True,
    auto_tail=False,
    head_pad=0.05,
    tail_pad=0.10,
    scene_threshold: float = SCENE_THRESHOLD,
) -> tuple[float, float]:
    dur = max(0.01, video_duration(webm))
    if not (auto_head or auto_tail):
        return 0.0, dur
    sc = detect_scene_changes(webm, scene_threshold)
    if not sc:
        return 0.0, dur
    start = max(0.0, sc[0] - head_pad) if auto_head else 0.0
    if auto_tail:
        end = min(dur, sc[-1] + tail_pad)
        take = max(0.01, end - start)
    else:
        take = max(0.01, dur - start)
    return start, take


def webm_to_mp4_precise(
    webm_path: Path,
    mp4_path: Path,
    start: float,
    take: float,
    out_w: int,
    out_h: int,
    fps: int,
):
    # Usa -ss DEPOIS do -i para seek preciso; mantém a mesma dimensão (sem barras).
    vf_chain = f"fps={fps}"
    cmd = [
        "ffmpeg",
        "-y",
        "-i",
        str(webm_path),
        "-ss",
        f"{start}",
        "-t",
        f"{take}",
        "-vf",
        vf_chain,
        "-c:v",
        "libx264",
        "-pix_fmt",
        PIX_FMT,
        "-profile:v",
        PROFILE,
        "-level:v",
        LEVEL,
        "-preset",
        PRESET,
        "-crf",
        str(CRF),
        "-movflags",
        "+faststart",
        "-an",
        str(mp4_path),
    ]
    subprocess.run(cmd, check=True)


def write_temp_html(tmpdir: Path, *, html_text: str = None, html_file=None) -> Path:
    """
    Grava o HTML recebido (texto ou arquivo enviado) no tmpdir e retorna o caminho.
    Mantém o nome base do upload (caso tenha assets relativos ao lado).
    """
    if html_text is not None:
        p = tmpdir / "index.html"
        p.write_text(html_text, encoding="utf-8")
        return p
    else:
        if html_file is None:
            raise ValueError("Nenhum HTML fornecido.")
        filename = getattr(html_file, "filename", None) or "upload.html"
        p = tmpdir / Path(filename).name
        html_file.save(p)
        return p


@app.route("/healthz")
def health():
    return {"ok": True}


def _needs_modern_renderer(html_path: Path) -> bool:
    """Heurística: detecta uso de CSS moderno (grid/flex) e outras pistas.
    Se encontrar, priorizamos Chromium para fidelidade.
    """
    try:
        txt = html_path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return False
    txt_low = txt.lower()
    patterns = (
        "display:grid",
        "grid-template",
        "grid-column",
        "grid-row",
        "@supports (display: grid)",
        "display:flex",
        "flex-wrap",
        "clamp(",
    )
    return any(p in txt_low for p in patterns)


@app.route('/html-para-imagem', methods=['POST'])
def html_para_imagem():
    """
    Converte um HTML em imagem.
    Aceita:
      - multipart/form-data com campo 'file' (arquivo .html)
      - OU multipart/form-data com campo 'html' (texto bruto do HTML)

    Parâmetros opcionais (form-data):
      - format: png|jpeg|jpg|webp (padrão: png)
      - width, height: dimensões em px (opcionais)
      - quality: qualidade para JPEG/WEBP (1-100)
      - transparent: 1/true para fundo transparente (apenas PNG)
      - filename: nome do arquivo de saída (opcional)

    Resposta: imagem (attachment)
    """
    html_file = request.files.get('file')
    html_text = request.form.get('html')
    if not html_file and not html_text:
        return jsonify({"error": "Envie 'file' (multipart) OU 'html' (texto)."}), 400

    fmt = (request.form.get('format') or request.form.get('formato') or 'png').lower()
    if fmt == 'jpg':
        fmt = 'jpeg'
    if fmt not in ('png', 'jpeg', 'webp'):
        return jsonify({"error": "Formato inválido. Use png, jpeg ou webp."}), 400

    renderer = (request.form.get('renderer') or request.form.get('renderizador') or 'auto').lower()
    if renderer not in ('wkhtml', 'chromium', 'playwright', 'chrome', 'auto'):
        renderer = 'auto'
    # Seleção de renderer:
    # - 'auto': usa wkhtml se disponível; senão, Chromium/Playwright
    # - 'wkhtml': se ausente, cai para Chromium
    if renderer == 'auto':
        # Prefira stack wkhtml se QUALQUER um (wkhtmltoimage OU wkhtmltopdf) estiver disponível
        if has_wkhtmltoimage() or has_wkhtmltopdf():
            renderer = 'wkhtml'
        else:
            renderer = 'chromium'
    elif renderer == 'wkhtml':
        # Fica em 'wkhtml' se houver wkhtmltoimage ou wkhtmltopdf; senão cai para chromium
        if not (has_wkhtmltoimage() or has_wkhtmltopdf()):
            renderer = 'chromium'

    # Dimensões
    width = request.form.get('width') or request.form.get('largura')
    height = request.form.get('height') or request.form.get('altura')
    try:
        width = int(width) if width else None
        height = int(height) if height else None
    except ValueError:
        return jsonify({"error": "width/height inválidos"}), 400

    # Qualidade (aplica-se a JPEG/WEBP)
    quality = request.form.get('quality') or request.form.get('qualidade')
    try:
        quality = int(quality) if quality else None
        if quality is not None and not (1 <= quality <= 100):
            raise ValueError
    except ValueError:
        return jsonify({"error": "quality deve ser inteiro entre 1 e 100"}), 400

    transparent = (
        str(request.form.get('transparent', '0')).lower() in ('1', 'true', 't', 'yes', 'y')
    )

    desired_name = request.form.get('filename') or request.form.get('nome_arquivo')
    if desired_name:
        desired_name = secure_filename(desired_name)
        # Garante extensão
        ext = f'.{fmt}'
        if not desired_name.lower().endswith(ext):
            desired_name = f"{desired_name}{ext}"
        if not desired_name or desired_name == ext:
            desired_name = None

    # Diretório temporário com limpeza após resposta (compatível com Windows)
    tmpdir_obj = tempfile.TemporaryDirectory(prefix="html2img_")
    tmpdir = Path(tmpdir_obj.name)

    @after_this_request
    def _cleanup(response):
        for _ in range(10):
            try:
                tmpdir_obj.cleanup()
                break
            except PermissionError:
                time.sleep(0.3)
        return response

    try:
        html_path = write_temp_html(tmpdir, html_text=html_text, html_file=html_file)
        out_path = tmpdir / f"output.{fmt}"

        # Se o HTML usa Grid/Flex, priorize Chromium para respeitar layout moderno
        if renderer in ('auto', 'wkhtml') and _needs_modern_renderer(html_path):
            renderer = 'chromium'

        if renderer in ('chromium', 'playwright', 'chrome'):
            # Renderização com Playwright para melhor fidelidade ao CSS
            try:
                dpr = request.form.get('dpr') or request.form.get('scale')
                try:
                    dpr = float(dpr) if dpr else 1.0
                except ValueError:
                    return jsonify({"error": "dpr/scale inválido"}), 400
                dpr = max(0.5, min(3.0, dpr))

                # Opções avançadas para fidelidade/layout
                media = request.form.get('media')
                if media:
                    media = media.lower()
                    if media not in ("screen", "print"):
                        media = None
                full_page = str(request.form.get('full_page', '0')).lower() in ("1", "true", "t", "yes", "y")
                # Se altura não foi informada, captura página inteira para evitar corte vertical
                if height is None:
                    full_page = True
                # CSS extra opcional para estabilizar layout/evitar sobreposição
                safe_layout = str(request.form.get('safe_layout', '0')).lower() in ("1", "true", "t", "yes", "y")
                extra_css = request.form.get('css')

                # espera extra opcional
                wait_ms = request.form.get('wait_ms') or request.form.get('espera_ms')
                try:
                    wait_ms = int(wait_ms) if wait_ms else 0
                except ValueError:
                    return jsonify({"error": "wait_ms inválido"}), 400

                # CSS de fallback de fontes (Inter -> Roboto/DejaVu/Liberation/Arial)
                font_fallback_css = (
                    """
                    body, * {
                      font-family: Inter, Roboto, 'DejaVu Sans', 'Liberation Sans', Arial, sans-serif !important;
                    }
                    """
                )

                base_png = chromium_screenshot(
                    html_path,
                    width=width,
                    height=height,
                    transparent=transparent,
                    dpr=dpr,
                    full_page=full_page,
                    media=media,
                    css_inject=font_fallback_css,
                )
                # Injeta CSS após carregamento para ajustes finos (quando solicitado)
                if safe_layout or extra_css:
                    with sync_playwright() as p:
                        try:
                            browser = p.chromium.launch(headless=True)
                        except Exception:
                            browser = p.chromium.launch(headless=True, args=['--no-sandbox'])
                        ctx = browser.new_context(viewport={"width": width or 1080, "height": height or 1350}, device_scale_factor=dpr)
                        page = ctx.new_page()
                        page.goto(file_url(html_path), wait_until="load")
                        prepare_page(page, zero_anim_delay=True)
                        if media in ("screen", "print"):
                            try:
                                page.emulate_media(media=media)
                            except Exception:
                                pass
                        css_payload = ""
                        if safe_layout:
                            css_payload += """
                            * { animation: none !important; transition: none !important; }
                            img { display:block; }
                            .hero { padding-bottom: calc(var(--hero-pad) + 16px) !important; }
                            .legend { inset: auto 10px 0 auto !important; }
                            """
                        if extra_css:
                            css_payload += str(extra_css)
                        if css_payload.strip():
                            page.add_style_tag(content=css_payload)
                        # Re-tira o screenshot com CSS aplicado
                        base_png = Path(tempfile.mkstemp(suffix=".png")[1])
                        page.screenshot(path=str(base_png), full_page=full_page, omit_background=bool(transparent))
                        ctx.close(); browser.close()
                if wait_ms:
                    time.sleep(wait_ms/1000.0)

                # Converter PNG base para formato desejado (se necessário)
                img = Image.open(str(base_png))
                save_params = {}
                if fmt == 'jpeg':
                    img = img.convert('RGB')
                    if quality:
                        save_params['quality'] = quality
                        save_params['optimize'] = True
                        save_params['progressive'] = True
                elif fmt == 'webp':
                    # Preserva alfa se existir
                    if quality:
                        save_params['quality'] = quality
                img.save(str(out_path), format=fmt.upper(), **save_params)
            except Exception as e:
                return jsonify({"error": f"Falha ao renderizar via Chromium/Playwright: {str(e)}"}), 500
        else:
            # Renderização via wkhtml stack (imgkit preferencial; senão wkhtmltopdf + rasterização)
            if has_wkhtmltoimage():
                # wkhtmltoimage não suporta algumas flags de wkhtmltopdf (ex.: --print-media-type)
                # Para máxima compatibilidade, renderizamos em PNG e depois convertemos ao formato final via Pillow.
                wk_options = {
                    "encoding": "utf-8",
                    "enable-local-file-access": None,
                    # dar tempo para webfonts/JS
                    "javascript-delay": 500,  # ms
                    "no-stop-slow-scripts": None,
                }
                if width:
                    wk_options["width"] = width
                    wk_options["disable-smart-width"] = None
                if height:
                    wk_options["height"] = height
                # Forçar saída base PNG do wkhtmltoimage
                wk_options["format"] = "png"
                base_png = tmpdir / "_wk_base.png"
                if transparent:
                    wk_options["transparent"] = None

                try:
                    imgkit.from_file(str(html_path), str(base_png), options=wk_options)
                except OSError as e:
                    msg = str(e)
                    return jsonify({"error": f"Falha no wkhtmltoimage: {msg}"}), 500

                # Converte PNG base para formato final (png/jpeg/webp)
                img = Image.open(str(base_png))
                save_params = {}
                if fmt == 'jpeg':
                    img = img.convert('RGB')
                    if quality:
                        save_params['quality'] = quality
                        save_params['optimize'] = True
                        save_params['progressive'] = True
                elif fmt == 'webp':
                    if quality:
                        save_params['quality'] = quality
                img.save(str(out_path), format=fmt.upper(), **save_params)
            elif has_wkhtmltopdf():
                # Fallback: gerar PDF e rasterizar primeira página
                pdf_path = Path(tempfile.mkstemp(suffix=".pdf")[1])
                try:
                    # Monta comando wkhtmltopdf
                    cmd = [
                        "wkhtmltopdf",
                        "--enable-local-file-access",
                        "--print-media-type",
                        str(html_path),
                        str(pdf_path),
                    ]
                    subprocess.run(cmd, check=True)
                except Exception as e:
                    return jsonify({"error": f"Falha no wkhtmltopdf: {str(e)}"}), 500

                try:
                    doc = fitz.open(str(pdf_path))
                    if len(doc) == 0:
                        return jsonify({"error": "PDF gerado está vazio"}), 500
                    page = doc[0]
                    # renderizar em boa resolução; ajusta conforme necessário
                    pix = page.get_pixmap(dpi=200)
                    img = Image.open(BytesIO(pix.tobytes("png")))
                    save_params = {}
                    if fmt == 'jpeg':
                        img = img.convert('RGB')
                        if quality:
                            save_params['quality'] = quality
                            save_params['optimize'] = True
                            save_params['progressive'] = True
                    elif fmt == 'webp':
                        if quality:
                            save_params['quality'] = quality
                    img.save(str(out_path), format=fmt.upper(), **save_params)
                    doc.close()
                except Exception as e:
                    return jsonify({"error": f"Falha ao rasterizar PDF: {str(e)}"}), 500
                finally:
                    try:
                        os.remove(pdf_path)
                    except Exception:
                        pass
            else:
                # Sem wkhtmltoimage e sem wkhtmltopdf: como último recurso, Chromium
                try:
                    dpr = request.form.get('dpr') or request.form.get('scale')
                    dpr = float(dpr) if dpr else 1.0
                    dpr = max(0.5, min(3.0, dpr))
                    font_fallback_css = (
                        """
                        body, * {
                          font-family: Inter, Roboto, 'DejaVu Sans', 'Liberation Sans', Arial, sans-serif !important;
                        }
                        """
                    )
                    base_png = chromium_screenshot(
                        html_path,
                        width=width,
                        height=height,
                        transparent=transparent,
                        dpr=dpr,
                        out_dir=tmpdir,
                        css_inject=font_fallback_css,
                    )
                    img = Image.open(str(base_png))
                    save_params = {}
                    if fmt == 'jpeg':
                        img = img.convert('RGB')
                        if quality:
                            save_params['quality'] = quality
                            save_params['optimize'] = True
                            save_params['progressive'] = True
                    elif fmt == 'webp':
                        if quality:
                            save_params['quality'] = quality
                    img.save(str(out_path), format=fmt.upper(), **save_params)
                except Exception as e2:
                    return jsonify({
                        "error": "Nenhum renderizador disponível. Instale 'wkhtmltopdf' (recomendado para Docker) ou os browsers do Playwright (python -m playwright install).",
                        "detail": str(e2)
                    }), 500

        # Pós-processamento para ajustar tamanho final sem distorção, se solicitado
        out_w = request.form.get('out_width') or request.form.get('largura_final')
        out_h = request.form.get('out_height') or request.form.get('altura_final')
        fit_mode = (request.form.get('fit') or request.form.get('ajuste') or 'contain').lower()
        bg_hex = (request.form.get('bg') or request.form.get('background') or 'ffffff').lstrip('#')

        def _parse_hex(h):
            try:
                if len(h) == 6:
                    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
            except Exception:
                pass
            return (255, 255, 255)

        if out_w or out_h:
            try:
                img = Image.open(str(out_path))
                src_w, src_h = img.size
                tw = int(out_w) if out_w else None
                th = int(out_h) if out_h else None

                # Se apenas um informado, mantemos proporção e não adicionamos padding por padrão
                if tw and not th:
                    scale = tw / src_w
                    nw, nh = max(1, int(round(src_w * scale))), max(1, int(round(src_h * scale)))
                    img = img.resize((nw, nh), Image.LANCZOS)
                elif th and not tw:
                    scale = th / src_h
                    nw, nh = max(1, int(round(src_w * scale))), max(1, int(round(src_h * scale)))
                    img = img.resize((nw, nh), Image.LANCZOS)
                else:
                    # Ambos informados: aplicar modo de ajuste
                    tw = max(1, int(tw or src_w))
                    th = max(1, int(th or src_h))
                    if fit_mode not in ('contain', 'cover', 'fill'):
                        fit_mode = 'contain'

                    if fit_mode == 'fill':
                        img = img.resize((tw, th), Image.LANCZOS)
                    else:
                        if fit_mode == 'contain':
                            scale = min(tw / src_w, th / src_h)
                        else:  # cover
                            scale = max(tw / src_w, th / src_h)
                        iw, ih = max(1, int(round(src_w * scale))), max(1, int(round(src_h * scale)))
                        img = img.resize((iw, ih), Image.LANCZOS)

                        if fit_mode == 'contain':
                            # canvas com padding
                            if fmt in ('png', 'webp') and transparent:
                                canvas = Image.new('RGBA', (tw, th), (0, 0, 0, 0))
                                if img.mode != 'RGBA':
                                    img = img.convert('RGBA')
                            else:
                                color = _parse_hex(bg_hex)
                                canvas = Image.new('RGB', (tw, th), color)
                                if img.mode not in ('RGB', 'RGBA'):
                                    img = img.convert('RGB')
                            ox = (tw - img.size[0]) // 2
                            oy = (th - img.size[1]) // 2
                            canvas.paste(img, (ox, oy), img if canvas.mode == 'RGBA' and img.mode == 'RGBA' else None)
                            img = canvas
                        else:  # cover: recortar centro
                            left = max(0, (img.size[0] - tw) // 2)
                            top = max(0, (img.size[1] - th) // 2)
                            img = img.crop((left, top, left + tw, top + th))

                # Salvar imagem final no mesmo caminho, preservando formato
                save_params = {}
                if fmt in ('jpeg',) and quality:
                    save_params['quality'] = quality
                    save_params['optimize'] = True
                    save_params['progressive'] = True
                if fmt == 'webp':
                    if quality:
                        save_params['quality'] = quality
                    # preservar alpha para webp se houver
                    if img.mode == 'RGBA':
                        save_params['lossless'] = False
                img.save(str(out_path), format=fmt.upper(), **save_params)
            except Exception as e:
                return jsonify({"error": f"Falha ao ajustar tamanho final: {str(e)}"}), 500

        mimetype = {
            'png': 'image/png',
            'jpeg': 'image/jpeg',
            'webp': 'image/webp',
        }[fmt]

        return send_file(
            str(out_path),
            mimetype=mimetype,
            as_attachment=True,
            download_name=(desired_name or f"convertido.{fmt}"),
        )
    except Exception as e:
        return jsonify({"error": f"Falha ao converter HTML em imagem: {str(e)}"}), 500


@app.route('/html-para-pdf', methods=['POST'])
def html_para_pdf():
    """
    Converte um HTML recebido em PDF.
    Aceita:
      - multipart/form-data com campo 'file' (arquivo .html)
      - OU multipart/form-data com campo 'html' (texto bruto do HTML)

    Parâmetros opcionais (form-data):
      - format: tamanho da página (ex.: A4, Letter). Padrão: A4
      - margin_top, margin_right, margin_bottom, margin_left (em mm, px, etc. wkhtmltopdf)

    Resposta: application/pdf (attachment)
    """
    html_file = request.files.get('file')
    html_text = request.form.get('html')
    if not html_file and not html_text:
        return jsonify({"error": "Envie 'file' (multipart) OU 'html' (texto)."}), 400

    # Parâmetros opcionais de página/margem
    page_format = request.form.get('format', 'A4')
    m_top = request.form.get('margin_top')
    m_right = request.form.get('margin_right')
    m_bottom = request.form.get('margin_bottom')
    m_left = request.form.get('margin_left')
    # nome do arquivo desejado
    desired_name = request.form.get('filename') or request.form.get('nome_arquivo')
    if desired_name:
        desired_name = secure_filename(desired_name)
        if not desired_name.lower().endswith('.pdf'):
            desired_name = f"{desired_name}.pdf"
        if not desired_name or desired_name == '.pdf':
            desired_name = None

    # Diretório temporário com limpeza pós-resposta (compatível com Windows)
    tmpdir_obj = tempfile.TemporaryDirectory(prefix="html2pdf_")
    tmpdir = Path(tmpdir_obj.name)

    @after_this_request
    def _cleanup(response):
        for _ in range(10):
            try:
                tmpdir_obj.cleanup()
                break
            except PermissionError:
                time.sleep(0.3)
        return response

    try:
        html_path = write_temp_html(tmpdir, html_text=html_text, html_file=html_file)
        pdf_path = tmpdir / "output.pdf"

        # Monta comando wkhtmltopdf
        cmd = [
            "wkhtmltopdf",
            "--enable-local-file-access",
            "--print-media-type",
        ]
        if page_format:
            cmd += ["-s", str(page_format)]
        if m_top:
            cmd += ["-T", str(m_top)]
        if m_right:
            cmd += ["-R", str(m_right)]
        if m_bottom:
            cmd += ["-B", str(m_bottom)]
        if m_left:
            cmd += ["-L", str(m_left)]

        cmd += [str(html_path), str(pdf_path)]

        try:
            subprocess.run(cmd, check=True)
        except FileNotFoundError:
            return jsonify({
                "error": "wkhtmltopdf não encontrado. Instale-o no sistema ou use a imagem Docker fornecida."
            }), 500

        return send_file(
            str(pdf_path),
            mimetype="application/pdf",
            as_attachment=True,
            download_name=(desired_name or "convertido.pdf"),
        )
    except Exception as e:
        return jsonify({"error": f"Falha ao converter HTML em PDF: {str(e)}"}), 500


@app.route("/render", methods=["POST"])
def render():
    """
    POST /render
    Form-data:
      - file: arquivo .html  (opcional se mandar 'html')
      - html: texto bruto do HTML (opcional se mandar 'file')
      - content_seconds (float, opcional) | default CONTENT_SECONDS
      - width,height (int, opcionais) | se omitidos e AUTO_SIZE_BODY=True → usa tamanho do <body>
      - target_fps (int, opcional) | default TARGET_FPS
      - auto_trim_head (bool), auto_trim_tail (bool) [0/1, true/false]
      - zero_anim_delay (bool)
      - scene_threshold, head_pad, tail_pad (opcionais)
    Resposta: video/mp4 (attachment)
    """
    try:
        ensure_ffmpeg()
    except Exception as e:
        return jsonify({"error": str(e)}), 500

    html_file = request.files.get("file")
    html_text = request.form.get("html")
    if not html_file and not html_text:
        return jsonify({"error": "Envie 'file' (multipart) OU 'html' (texto)."}), 400

    # parâmetros opcionais
    try:
        content_seconds = float(request.form.get("content_seconds", CONTENT_SECONDS))
    except ValueError:
        return jsonify({"error": "content_seconds inválido"}), 400
    try:
        target_fps = int(request.form.get("target_fps", TARGET_FPS))
    except ValueError:
        return jsonify({"error": "target_fps inválido"}), 400

    auto_trim_head = (
        str(request.form.get("auto_trim_head", str(AUTO_TRIM_HEAD))).lower()
        in ("1", "true", "t", "yes", "y")
    )
    auto_trim_tail = (
        str(request.form.get("auto_trim_tail", str(AUTO_TRIM_TAIL))).lower()
        in ("1", "true", "t", "yes", "y")
    )
    zero_anim_delay = (
        str(request.form.get("zero_anim_delay", str(ZERO_ANIM_DELAY))).lower()
        in ("1", "true", "t", "yes", "y")
    )

    scene_threshold = request.form.get("scene_threshold")
    head_pad = request.form.get("head_pad")
    tail_pad = request.form.get("tail_pad")
    try:
        scene_threshold = (
            float(scene_threshold) if scene_threshold is not None else SCENE_THRESHOLD
        )
        head_pad = float(head_pad) if head_pad is not None else HEAD_PAD_S
        tail_pad = float(tail_pad) if tail_pad is not None else TAIL_PAD_S
    except ValueError:
        return jsonify({"error": "scene_threshold/head_pad/tail_pad inválidos"}), 400

    width = request.form.get("width")
    height = request.form.get("height")
    try:
        width = int(width) if width else None
        height = int(height) if height else None
    except ValueError:
        return jsonify({"error": "width/height inválidos"}), 400

    # -------- NÃO usar TemporaryDirectory como context manager (Windows lock) --------
    tmpdir_obj = tempfile.TemporaryDirectory(prefix="html2mp4_")
    tmpdir = Path(tmpdir_obj.name)

    @after_this_request
    def _cleanup(response):
        # Tenta limpar após o envio (Windows pode segurar handle por um pouco)
        for _ in range(10):
            try:
                tmpdir_obj.cleanup()
                break
            except PermissionError:
                time.sleep(0.3)
        return response

    try:
        html_path = write_temp_html(
            tmpdir, html_text=html_text, html_file=html_file
        )

        # 1) descobrir tamanho
        if AUTO_SIZE_BODY and (width is None or height is None):
            w, h = measure_body_size(html_path)
        else:
            w = width or 1080
            h = height or 1350

        # 2) gravar WEBM com buffers
        total = content_seconds + BUFFER_HEAD_S + BUFFER_TAIL_S
        webm = record_webm(
            html_path, total, w, h, tmpdir, zero_anim_delay=zero_anim_delay
        )

        # 3) auto-trim do início (e opcional do final), SEM cortar conteúdo
        start, take = compute_trim(
            webm,
            auto_head=auto_trim_head,
            auto_tail=auto_trim_tail,
            head_pad=head_pad,
            tail_pad=tail_pad,
            scene_threshold=scene_threshold,
        )

        # 4) converter para MP4 com seek preciso (sem barras)
        mp4_path = tmpdir / "output.mp4"
        webm_to_mp4_precise(webm, mp4_path, start, take, w, h, target_fps)

        # 5) devolver (passe string; Flask abre/fecha o arquivo)
        resp = send_file(
            str(mp4_path),
            mimetype="video/mp4",
            as_attachment=True,
            download_name="render.mp4",
            max_age=0,
            conditional=True,
        )
        # Evita cache agressivo de proxies/browsers
        resp.headers["Cache-Control"] = "no-store, no-cache, must-revalidate, max-age=0"
        resp.headers["Pragma"] = "no-cache"
        return resp
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/pdf-para-docx', methods=['POST'])
def pdf_para_docx():
    """
    Converte um PDF em um arquivo .docx com opção EDITÁVEL.
    - Modo 'editable' (padrão): usa pdf2docx para gerar DOCX com textos/caixas editáveis.
    - Modo 'raster': insere cada página como imagem (fidelidade visual máxima, não editável).

    Entrada (multipart/form-data):
      - file: PDF (alternativa: pdf_url)
      - pdf_url: URL para um PDF (opcional se enviar file)
      - mode: editable | raster (padrão: editable)
      - dpi: (raster) resolução para rasterização (padrão: 300)
      - page_size: (raster) A4 | Letter (padrão: A4)
      - margin_top, margin_right, margin_bottom, margin_left (em mm, opcionais — raster)
      - start_page, end_page: intervalo (0-based inclusive/exclusive) para conversão (editable)
      - filename: nome do arquivo de saída (opcional)
    """
    # Fonte: arquivo ou URL
    uploaded = request.files.get('file')
    pdf_url = request.form.get('pdf_url')
    if not uploaded and not pdf_url:
        return jsonify({"error": "Envie 'file' (PDF) ou 'pdf_url'."}), 400

    # Nome desejado
    desired_name = request.form.get('filename') or request.form.get('nome_arquivo')
    if desired_name:
        desired_name = secure_filename(desired_name)
        if not desired_name.lower().endswith('.docx'):
            desired_name = f"{desired_name}.docx"

    # Parâmetros
    mode = (request.form.get('mode') or 'editable').strip().lower()
    # "exact" é um alias para raster (fidelidade 100% visual)
    exact = str(request.form.get('exact', '0')).lower() in ('1','true','t','yes','y')
    if exact:
        mode = 'raster'
    if mode not in ('editable', 'raster', 'exact', 'hybrid'):
        mode = 'editable'

    try:
        dpi = int(request.form.get('dpi', '300'))
        if dpi < 72 or dpi > 600:
            dpi = 300
    except ValueError:
        dpi = 300

    page_size = (request.form.get('page_size') or 'A4').strip().lower()
    if page_size not in ('a4', 'letter'):
        page_size = 'a4'

    # Margens em mm (se Mm disponível)
    def _mm(val, default):
        if Mm is None or val is None:
            return default
        try:
            return Mm(float(val))
        except Exception:
            return default

    margin_top = request.form.get('margin_top')
    margin_right = request.form.get('margin_right')
    margin_bottom = request.form.get('margin_bottom')
    margin_left = request.form.get('margin_left')

    # Obter bytes do PDF
    pdf_bytes = None
    tmp_pdf_path = None
    try:
        if uploaded:
            # Ler stream para bytes
            uploaded.stream.seek(0)
            pdf_bytes = uploaded.read()
        else:
            resp = requests.get(pdf_url, timeout=30)
            resp.raise_for_status()
            pdf_bytes = resp.content
        if not pdf_bytes:
            return jsonify({"error": "PDF vazio."}), 400

        # Caminho PDF temporário para engines que exigem arquivo
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
            tmp_pdf.write(pdf_bytes)
            tmp_pdf_path = tmp_pdf.name

        out_path = tempfile.NamedTemporaryFile(delete=False, suffix='.docx').name

        # Intervalo opcional de páginas (0-based). Se não informado, converte todas.
        start_page = request.form.get('start_page')
        end_page = request.form.get('end_page')
        try:
            start_i = int(start_page) if start_page is not None else 0
            if start_i < 0:
                start_i = 0
        except ValueError:
            start_i = 0
        try:
            end_i = int(end_page) if end_page is not None else None
            if end_i is not None and end_i < 0:
                end_i = None
        except ValueError:
            end_i = None

        if mode == 'editable':
            if Converter is None:
                return jsonify({"error": "Conversão editável requer 'pdf2docx'. Instale as dependências (requirements/Docker)."}), 500
            # Conversão via pdf2docx
            cv = Converter(tmp_pdf_path)
            cv.convert(out_path, start=start_i, end=end_i)
            cv.close()
        elif mode == 'hybrid':
            # Híbrido: insere imagem da página e adiciona texto extraído embaixo (editável)
            doc_pdf = fitz.open(stream=pdf_bytes, filetype='pdf')
            if len(doc_pdf) == 0:
                return jsonify({"error": "PDF sem páginas."}), 400
            if Document is None:
                return jsonify({"error": "python-docx não está disponível no ambiente."}), 500
            docx = Document()
            section = docx.sections[0]
            if Mm:
                if page_size == 'a4':
                    section.page_width = Mm(210)
                    section.page_height = Mm(297)
                else:
                    section.page_width = Mm(215.9)
                    section.page_height = Mm(279.4)
                section.top_margin = _mm(margin_top, section.top_margin)
                section.right_margin = _mm(margin_right, section.right_margin)
                section.bottom_margin = _mm(margin_bottom, section.bottom_margin)
                section.left_margin = _mm(margin_left, section.left_margin)
            page_width_emu = section.page_width - section.left_margin - section.right_margin if Mm else None
            tmpdir_obj = tempfile.TemporaryDirectory(prefix="pdf2docx_hybrid_")
            tmpdir = Path(tmpdir_obj.name)
            total_pages = len(doc_pdf)
            real_end = end_i if (end_i is not None and end_i <= total_pages) else total_pages
            real_start = max(0, min(start_i, total_pages))
            for idx in range(real_start, real_end):
                page = doc_pdf[idx]
                # Imagem da página
                pix = page.get_pixmap(dpi=dpi, alpha=False)
                img_path = tmpdir / f"page_{idx+1}.png"
                img_path.write_bytes(pix.tobytes('png'))
                if page_width_emu is not None:
                    pic = docx.add_picture(str(img_path))
                    pic.width = page_width_emu
                else:
                    docx.add_picture(str(img_path))
                # Texto extraído (editável) logo abaixo
                try:
                    text = page.get_text("text") or ""
                except Exception:
                    text = ""
                if text.strip():
                    for para in text.splitlines():
                        p = docx.add_paragraph(para)
                # Quebra de página se não for a última
                if idx < (real_end - 1):
                    docx.add_page_break()
            docx.save(out_path)
        else:
            # Raster (fidelidade visual), docx com imagens
            doc_pdf = fitz.open(stream=pdf_bytes, filetype='pdf')
            if len(doc_pdf) == 0:
                return jsonify({"error": "PDF sem páginas."}), 400
            if Document is None:
                return jsonify({"error": "python-docx não está disponível no ambiente."}), 500
            docx = Document()
            section = docx.sections[0]
            if Mm:
                if page_size == 'a4':
                    section.page_width = Mm(210)
                    section.page_height = Mm(297)
                else:
                    section.page_width = Mm(215.9)
                    section.page_height = Mm(279.4)
                section.top_margin = _mm(margin_top, section.top_margin)
                section.right_margin = _mm(margin_right, section.right_margin)
                section.bottom_margin = _mm(margin_bottom, section.bottom_margin)
                section.left_margin = _mm(margin_left, section.left_margin)
            page_width_emu = section.page_width - section.left_margin - section.right_margin if Mm else None
            tmpdir_obj = tempfile.TemporaryDirectory(prefix="pdf2docx_")
            tmpdir = Path(tmpdir_obj.name)
            total_pages = len(doc_pdf)
            real_end = end_i if (end_i is not None and end_i <= total_pages) else total_pages
            real_start = max(0, min(start_i, total_pages))
            for idx in range(real_start, real_end):
                page = doc_pdf[idx]
                pix = page.get_pixmap(dpi=dpi, alpha=False)
                img_path = tmpdir / f"page_{idx+1}.png"
                img_path.write_bytes(pix.tobytes('png'))
                if page_width_emu is not None:
                    pic = docx.add_picture(str(img_path))
                    pic.width = page_width_emu
                else:
                    docx.add_picture(str(img_path))
                if idx < (real_end - 1):
                    docx.add_page_break()
            docx.save(out_path)

        # Enviar
        return send_file(
            out_path,
            mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            as_attachment=True,
            download_name=(desired_name or 'convertido.docx')
        )
    except requests.RequestException as e:
        return jsonify({"error": f"Falha ao baixar PDF: {str(e)}"}), 400
    except Exception as e:
        return jsonify({"error": f"Falha ao converter PDF para DOCX: {str(e)}"}), 500

if __name__ == '__main__':
    app.run(host="0.0.0.0", port=5000)
